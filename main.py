import sys
import os
import socket
import subprocess
from PyQt5.QtWidgets import (QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout,
                             QTextEdit, QLineEdit, QPushButton, QLabel, QListWidget, QCheckBox, 
                             QFrame, QMessageBox, QDialog, QProgressBar)
from PyQt5.QtCore import Qt, QThread, pyqtSignal
from pypdf import PdfReader
from pptx import Presentation
import ollama
from openai import OpenAI

# --- Internet Kontrolü ---
def internet_var_mi():
    try:
        socket.create_connection(("8.8.8.8", 53), timeout=2)
        return True
    except OSError:
        return False

# --- Arka Planda Canlı İndirme İşçisi ---
class ModelIndirmeIscisi(QThread):
    ilerleme = pyqtSignal(int)
    durum_mesaji = pyqtSignal(str)
    tamamlandi = pyqtSignal()
    hata = pyqtSignal(str)

    def run(self):
        try:
            for yanit in ollama.pull('qwen2.5:7b', stream=True):
                durum = yanit.get('status', '')
                self.durum_mesaji.emit(durum)
                
                if 'total' in yanit and 'completed' in yanit:
                    if yanit['total'] > 0:
                        yuzde = int((yanit['completed'] / yanit['total']) * 100)
                        self.ilerleme.emit(yuzde)
            
            self.ilerleme.emit(100)
            self.durum_mesaji.emit("✅ İndirme Tamamlandı!")
            self.tamamlandi.emit()
        except Exception as e:
            self.hata.emit(f"Hata: İndirme başarısız oldu. Ağ bağlantınızı kontrol edin.\nDetay: {str(e)}")

# --- Başlangıç Kontrol Kutusu (Sihirbaz) ---
class BaslangicSihirbazi(QDialog):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setWindowTitle("Sistem Kontrolü")
        self.setFixedSize(450, 200)
        self.setWindowFlags(Qt.Dialog | Qt.CustomizeWindowHint | Qt.WindowTitleHint) 
        self.setModal(True)
        self.setStyleSheet("""
            QDialog { background-color: #1a1d24; border: 2px solid #333; border-radius: 10px; }
            QLabel { color: white; font-size: 14px; font-family: 'Segoe UI'; }
            QProgressBar { border: 1px solid #333; border-radius: 5px; text-align: center; color: white; background-color: #0b0c10; height: 25px;}
            QProgressBar::chunk { background-color: #0078D7; width: 20px; }
            QPushButton { padding: 10px; border-radius: 5px; font-weight: bold; font-family: 'Segoe UI'; }
            QPushButton#BtnMavi { background-color: #0078D7; color: white; }
            QPushButton#BtnMavi:hover { background-color: #005a9e; }
            QPushButton#BtnYesil { background-color: #2e7d32; color: white; }
            QPushButton#BtnYesil:hover { background-color: #1b5e20; }
            QPushButton:disabled { background-color: #333; color: #777; }
        """)

        layout = QVBoxLayout(self)
        
        self.bilgi_etiketi = QLabel("Sistem donanımları ve yapay zeka modelleri taranıyor...")
        self.bilgi_etiketi.setWordWrap(True)
        self.bilgi_etiketi.setAlignment(Qt.AlignCenter)
        layout.addWidget(self.bilgi_etiketi)

        self.durum_etiketi = QLabel("")
        self.durum_etiketi.setStyleSheet("color: #aaa; font-style: italic; font-size: 12px;")
        self.durum_etiketi.setAlignment(Qt.AlignCenter)
        layout.addWidget(self.durum_etiketi)

        self.progress_bar = QProgressBar(self)
        self.progress_bar.setValue(0)
        self.progress_bar.hide() 
        layout.addWidget(self.progress_bar)

        self.btn_aksiyon = QPushButton("Bekleniyor...")
        self.btn_aksiyon.setObjectName("BtnMavi")
        layout.addWidget(self.btn_aksiyon)

        self.modeli_denetle()

    def modeli_denetle(self):
        try:
            sonuc = subprocess.run(["ollama", "list"], capture_output=True, text=True, check=True)
            if "qwen2.5:7b" in sonuc.stdout:
                self.bilgi_etiketi.setText("✅ Qwen 2.5 (Yerel Zeka) sistemi hazır ve aktif.")
                self.bilgi_etiketi.setStyleSheet("color: #4caf50; font-weight: bold; font-size: 15px;")
                self.btn_aksiyon.setText("Devam")
                self.btn_aksiyon.setObjectName("BtnYesil")
                self.btn_aksiyon.clicked.connect(self.accept) 
            else:
                self.bilgi_etiketi.setText("⚠️ Qwen 2.5 (4.7 GB) yerel modeli bulunamadı.\nÇevrimdışı çalışabilmek için modelin indirilmesi gerekiyor.")
                self.bilgi_etiketi.setStyleSheet("color: #ff9800; font-weight: bold;")
                self.btn_aksiyon.setText("İndir ve Kur")
                self.btn_aksiyon.setObjectName("BtnMavi")
                self.btn_aksiyon.clicked.connect(self.indirmeyi_baslat)
        except Exception:
            self.bilgi_etiketi.setText("❌ Sistemde 'Ollama' uygulaması bulunamadı!\nCihazınızda Ollama kurulu değilse sadece Çevrimiçi Mod çalışır.")
            self.bilgi_etiketi.setStyleSheet("color: #f44336; font-weight: bold;")
            self.btn_aksiyon.setText("Çevrimiçi Modla Devam Et")
            self.btn_aksiyon.setObjectName("BtnMavi")
            self.btn_aksiyon.clicked.connect(self.accept)

    def indirmeyi_baslat(self):
        self.btn_aksiyon.setEnabled(False)
        self.btn_aksiyon.setText("İndiriliyor, Uygulamayı Kapatmayın...")
        self.progress_bar.show()
        
        self.isci = ModelIndirmeIscisi()
        self.isci.ilerleme.connect(self.progress_bar.setValue)
        self.isci.durum_mesaji.connect(self.durum_etiketi.setText)
        self.isci.tamamlandi.connect(self.indirme_bitti)
        self.isci.hata.connect(self.hata_olustu)
        self.isci.start()

    def indirme_bitti(self):
        self.bilgi_etiketi.setText("✅ Kurulum Başarılı! Sistem kullanıma hazır.")
        self.bilgi_etiketi.setStyleSheet("color: #4caf50; font-weight: bold;")
        self.durum_etiketi.setText("")
        self.btn_aksiyon.setEnabled(True)
        self.btn_aksiyon.setText("Devam")
        self.btn_aksiyon.setObjectName("BtnYesil")
        self.btn_aksiyon.disconnect() 
        self.btn_aksiyon.clicked.connect(self.accept)

    def hata_olustu(self, hata_mesaji):
        self.btn_aksiyon.setEnabled(True)
        self.btn_aksiyon.setText("Yeniden Dene")
        self.durum_etiketi.setText(hata_mesaji)
        self.durum_etiketi.setStyleSheet("color: #f44336;")


# --- Sohbet İşçisi ---
class YapayZekaIscisi(QThread):
    yanit_geldi = pyqtSignal(str, str)

    def __init__(self, soru, metin, api_key, openai_zorla):
        super().__init__()
        self.soru = soru
        self.metin = metin
        self.api_key = api_key
        self.openai_zorla = openai_zorla

    def run(self):
        sinirli_metin = self.metin[-7000:] if len(self.metin) > 7000 else self.metin
        sistem_talimati = f"Sen cihaz üzerinde çalışan gizlilik odaklı kurumsal bir asistansın. Verilen BELGELER dışında bilgi kullanma.\n\nBELGELER:\n{sinirli_metin}"
        used_model = "Yerel Model (Qwen2.5:7b)"
        
        if self.openai_zorla and self.api_key and internet_var_mi():
            try:
                if self.api_key.startswith("gsk_"):
                    client = OpenAI(api_key=self.api_key, base_url="https://api.groq.com/openai/v1")
                    model_secimi = "llama3-8b-8192"
                    used_model = "Bulut Model (Groq Llama-3)"
                else:
                    client = OpenAI(api_key=self.api_key)
                    model_secimi = "gpt-3.5-turbo"
                    used_model = "Bulut Model (OpenAI GPT)"

                response = client.chat.completions.create(
                    model=model_secimi,
                    messages=[
                        {"role": "system", "content": sistem_talimati},
                        {"role": "user", "content": self.soru}
                    ],
                    temperature=0.1, max_tokens=500
                )
                self.yanit_geldi.emit(response.choices[0].message.content, used_model)
                return
            except Exception:
                used_model = "Yerel Model (Qwen2.5) - Fallback"

        try:
            cevap = ollama.chat(model='qwen2.5:7b', messages=[
                {'role': 'system', 'content': sistem_talimati},
                {'role': 'user', 'content': self.soru}
            ], options={'temperature': 0.1, 'num_predict': 300, 'num_ctx': 4096})
            self.yanit_geldi.emit(cevap['message']['content'], used_model)
        except Exception as e:
            self.yanit_geldi.emit(f"❌ Hata: Yerel model çalıştırılamadı.", "Hata")


# --- Ana Uygulama ---
class YerelAsistanApp(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("Yerel Yapay Zeka Belge Asistanı")
        self.resize(1100, 700)
        self.setAcceptDrops(True)
        self.belge_icerigi = ""
        self.karanlik_mod_aktif = True

        self.arayuzu_kur()
        self.temayi_guncelle()

    def showEvent(self, event):
        super().showEvent(event)
        self.baslangic_ekranini_goster()

    def baslangic_ekranini_goster(self):
        self.overlay = QWidget(self)
        self.overlay.setStyleSheet("background-color: rgba(0, 0, 0, 190);")
        self.overlay.resize(self.size())
        self.overlay.show()

        self.dialog = BaslangicSihirbazi(self)
        self.dialog.exec_()

        self.overlay.hide()

    def resizeEvent(self, event):
        super().resizeEvent(event)
        if hasattr(self, 'overlay') and self.overlay.isVisible():
            self.overlay.resize(self.size())

    def arayuzu_kur(self):
        merkez_widget = QWidget()
        self.setCentralWidget(merkez_widget)
        ana_layout = QHBoxLayout(merkez_widget)
        ana_layout.setContentsMargins(15, 15, 15, 15)
        ana_layout.setSpacing(15)

        self.sol_panel = QFrame()
        self.sol_panel.setObjectName("Panel")
        self.sol_panel.setFixedWidth(280)
        sol_layout = QVBoxLayout(self.sol_panel)

        belgeler_baslik = QLabel("Belgeler")
        belgeler_baslik.setStyleSheet("font-size: 16px; font-weight: bold; margin-bottom: 5px;")
        
        self.dosya_etiket = QLabel("📂 Dosya veya klasörü\nburaya sürükleyin\n\nTXT • PDF • PPTX")
        self.dosya_etiket.setObjectName("SurukleAlani")
        self.dosya_etiket.setAlignment(Qt.AlignCenter)
        self.dosya_etiket.setMinimumHeight(120)
        
        self.dosya_listesi = QListWidget()
        self.dosya_listesi.setObjectName("DosyaListesi")

        api_baslik = QLabel("🔑 Bulut API Anahtarı (Opsiyonel)")
        api_baslik.setStyleSheet("font-size: 13px; font-weight: bold; margin-top: 10px;")
        self.api_girdisi = QLineEdit()
        self.api_girdisi.setPlaceholderText("sk-... veya gsk_...")
        self.api_girdisi.setEchoMode(QLineEdit.Password) 

        self.btn_belge_temizle = QPushButton("🗑️ Belgeleri Temizle")
        self.btn_belge_temizle.setObjectName("BtnKirmizi")
        self.btn_belge_temizle.clicked.connect(self.belgeleri_temizle)

        self.btn_sohbet_temizle = QPushButton("🧠 Sohbet Hafızasını Temizle")
        self.btn_sohbet_temizle.setObjectName("BtnGri")
        self.btn_sohbet_temizle.clicked.connect(self.sohbeti_temizle)

        self.btn_tema_degistir = QPushButton("☀️ Açık Moda Geç")
        self.btn_tema_degistir.setObjectName("BtnGri")
        self.btn_tema_degistir.clicked.connect(self.tema_degistir)

        sol_layout.addWidget(belgeler_baslik)
        sol_layout.addWidget(self.dosya_etiket)
        sol_layout.addWidget(self.dosya_listesi)
        sol_layout.addWidget(api_baslik)
        sol_layout.addWidget(self.api_girdisi)
        sol_layout.addSpacing(10)
        sol_layout.addWidget(self.btn_belge_temizle)
        sol_layout.addWidget(self.btn_sohbet_temizle)
        sol_layout.addWidget(self.btn_tema_degistir)
        ana_layout.addWidget(self.sol_panel)

        self.sag_panel = QFrame()
        self.sag_panel.setObjectName("Panel")
        sag_layout = QVBoxLayout(self.sag_panel)

        ust_bar_layout = QHBoxLayout()
        baslik_sag = QLabel("Karma Yapay Zeka Belge Asistanı")
        baslik_sag.setStyleSheet("font-size: 16px; font-weight: bold;")
        self.durum_etiketi = QLabel("🟢 Hazır")
        self.durum_etiketi.setObjectName("DurumEtiketi")
        self.durum_etiketi.setAlignment(Qt.AlignRight)
        
        ust_bar_layout.addWidget(baslik_sag)
        ust_bar_layout.addWidget(self.durum_etiketi)

        self.sohbet_ekrani = QTextEdit()
        self.sohbet_ekrani.setReadOnly(True)
        self.sohbet_ekrani.setObjectName("SohbetEkrani")
        self.ilk_mesaji_yazdir()

        self.openai_check = QCheckBox("Çevrimiçi Modu (Bulut API) Kullan")
        
        girdi_layout = QHBoxLayout()
        self.soru_girdisi = QLineEdit()
        self.soru_girdisi.setPlaceholderText("Belgelerle ilgili sorunuzu yazın...")
        self.soru_girdisi.returnPressed.connect(self.soru_sor)
        
        self.btn_gonder = QPushButton("Gönder")
        self.btn_gonder.setObjectName("BtnMavi")
        self.btn_gonder.setFixedWidth(100)
        self.btn_gonder.clicked.connect(self.soru_sor)

        girdi_layout.addWidget(self.soru_girdisi)
        girdi_layout.addWidget(self.btn_gonder)

        sag_layout.addLayout(ust_bar_layout)
        sag_layout.addWidget(self.sohbet_ekrani)
        sag_layout.addWidget(self.openai_check)
        sag_layout.addLayout(girdi_layout)

        ana_layout.addWidget(self.sag_panel)

    def tema_degistir(self):
        self.karanlik_mod_aktif = not self.karanlik_mod_aktif
        self.btn_tema_degistir.setText("☀️ Açık Moda Geç" if self.karanlik_mod_aktif else "🌙 Karanlık Moda Geç")
        self.temayi_guncelle()

    def temayi_guncelle(self):
        if self.karanlik_mod_aktif:
            self.setStyleSheet("""
                QMainWindow { background-color: #0b0c10; }
                QFrame#Panel { background-color: #15171e; border-radius: 8px; border: 1px solid #222; }
                QLabel { color: #ffffff; font-family: 'Segoe UI', Arial; }
                QLabel#DurumEtiketi { color: #4caf50; font-weight: bold; }
                QLabel#SurukleAlani { background-color: #1a1d24; border: 2px dashed #333; border-radius: 8px; color: #888; font-size: 13px; }
                QListWidget#DosyaListesi { background-color: #0b0c10; border: 1px solid #222; border-radius: 6px; color: #ddd; padding: 5px; }
                QTextEdit#SohbetEkrani { background-color: #0f1115; border: 1px solid #222; border-radius: 6px; color: #ddd; padding: 15px; font-size: 14px; font-family: 'Segoe UI', Arial; }
                QCheckBox { color: #aaa; font-family: 'Segoe UI'; }
                QLineEdit { background-color: #0b0c10; border: 1px solid #333; border-radius: 6px; padding: 10px; color: white; font-size: 14px; }
                QLineEdit:focus { border: 1px solid #0078D7; }
                QPushButton { padding: 10px; border-radius: 6px; font-weight: bold; font-family: 'Segoe UI'; border: none; }
                QPushButton#BtnMavi { background-color: #0078D7; color: white; }
                QPushButton#BtnMavi:hover { background-color: #005a9e; }
                QPushButton#BtnKirmizi { background-color: #3b1f1f; color: #ff9999; border: 1px solid #552222; }
                QPushButton#BtnKirmizi:hover { background-color: #4a2727; }
                QPushButton#BtnGri { background-color: #22252b; color: #ccc; border: 1px solid #333; }
                QPushButton#BtnGri:hover { background-color: #2d3139; }
                QPushButton:disabled { background-color: #222; color: #555; }
            """)
        else:
            self.setStyleSheet("""
                QMainWindow { background-color: #eef0f4; }
                QFrame#Panel { background-color: #ffffff; border-radius: 8px; border: 1px solid #ccc; }
                QLabel { color: #111111; font-family: 'Segoe UI', Arial; }
                QLabel#DurumEtiketi { color: #2e7d32; font-weight: bold; }
                QLabel#SurukleAlani { background-color: #f8f9fa; border: 2px dashed #aaa; border-radius: 8px; color: #555; font-size: 13px; }
                QListWidget#DosyaListesi { background-color: #f4f5f7; border: 1px solid #ddd; border-radius: 6px; color: #333; padding: 5px; }
                QTextEdit#SohbetEkrani { background-color: #f8f9fa; border: 1px solid #ddd; border-radius: 6px; color: #222; padding: 15px; font-size: 14px; font-family: 'Segoe UI', Arial; }
                QCheckBox { color: #444; font-family: 'Segoe UI'; }
                QLineEdit { background-color: #ffffff; border: 1px solid #bbb; border-radius: 6px; padding: 10px; color: #111; font-size: 14px; }
                QLineEdit:focus { border: 1px solid #0078D7; }
                QPushButton { padding: 10px; border-radius: 6px; font-weight: bold; font-family: 'Segoe UI'; border: none; }
                QPushButton#BtnMavi { background-color: #0078D7; color: white; }
                QPushButton#BtnMavi:hover { background-color: #005a9e; }
                QPushButton#BtnKirmizi { background-color: #ffeaea; color: #c62828; border: 1px solid #ffcaca; }
                QPushButton#BtnKirmizi:hover { background-color: #ffdbdb; }
                QPushButton#BtnGri { background-color: #e0e4e8; color: #333; border: 1px solid #ccc; }
                QPushButton#BtnGri:hover { background-color: #d1d6dc; }
                QPushButton:disabled { background-color: #ddd; color: #999; }
            """)

    def ilk_mesaji_yazdir(self):
        self.sohbet_ekrani.append("🤖 Sistem hazır.\n\nSol panele belge sürükleyip bırakın.\nÇevrimiçi mod için sol tarafa API anahtarınızı (Groq veya OpenAI) girebilirsiniz.\n")

    def sohbeti_temizle(self):
        self.sohbet_ekrani.clear()
        self.ilk_mesaji_yazdir()

    def belgeleri_temizle(self):
        self.belge_icerigi = ""
        self.dosya_listesi.clear()
        self.sohbet_ekrani.append("🗑️ Sistemdeki tüm belgeler silindi.\n")

    def dragEnterEvent(self, event):
        if event.mimeData().hasUrls(): event.acceptProposedAction()

    def dropEvent(self, event):
        for url in event.mimeData().urls():
            yol = url.toLocalFile()
            if os.path.isdir(yol):
                for kok_dizin, _, dosyalar in os.walk(yol):
                    for dosya in dosyalar: self.dosyayi_oku_ve_ekle(os.path.join(kok_dizin, dosya))
            else: self.dosyayi_oku_ve_ekle(yol)

    def dosyayi_oku_ve_ekle(self, dosya_yolu):
        uzanti = dosya_yolu.lower().split('.')[-1]
        if uzanti not in ['txt', 'pdf', 'pptx']: return
        isim = os.path.basename(dosya_yolu)
        icerik = ""
        try:
            if uzanti == 'txt':
                with open(dosya_yolu, 'r', encoding='utf-8') as f: icerik = f.read()
            elif uzanti == 'pdf':
                reader = PdfReader(dosya_yolu)
                for sayfa in reader.pages:
                    metin = sayfa.extract_text()
                    if metin: icerik += metin + "\n"
            elif uzanti == 'pptx':
                prs = Presentation(dosya_yolu)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"): icerik += shape.text + "\n"
            self.belge_icerigi += f"\n--- {isim} İçeriği ---\n{icerik}\n"
            self.dosya_listesi.addItem(f"📄 {isim}")
        except: pass

    def soru_sor(self):
        soru = self.soru_girdisi.text().strip()
        if not soru or not self.belge_icerigi: return
        
        girilen_api = self.api_girdisi.text().strip()
        gpt_aktif = self.openai_check.isChecked()
        
        if gpt_aktif and not girilen_api:
            QMessageBox.warning(self, "API Anahtarı Eksik", "Çevrimiçi (Bulut) modunu kullanmak için sol panele geçerli bir API anahtarı girmelisiniz.")
            return

        self.sohbet_ekrani.append(f"👤 Sen: {soru}")
        self.soru_girdisi.clear()
        
        self.soru_girdisi.setEnabled(False)
        self.btn_gonder.setEnabled(False)
        self.durum_etiketi.setText("⏳ Düşünüyor...")
        self.durum_etiketi.setStyleSheet("color: #ff9800; font-weight: bold;")

        self.worker = YapayZekaIscisi(soru, self.belge_icerigi, girilen_api, gpt_aktif)
        self.worker.yanit_geldi.connect(self.cevap_yazdir)
        self.worker.start()

    def cevap_yazdir(self, cevap, kullanilan_model):
        self.sohbet_ekrani.append(f"🤖 Asistan ({kullanilan_model}):\n{cevap}\n{'-'*40}\n")
        self.durum_etiketi.setText("🟢 Hazır")
        if self.karanlik_mod_aktif:
            self.durum_etiketi.setStyleSheet("color: #4caf50; font-weight: bold;")
        else:
            self.durum_etiketi.setStyleSheet("color: #2e7d32; font-weight: bold;")
        
        self.soru_girdisi.setEnabled(True)
        self.btn_gonder.setEnabled(True)
        self.soru_girdisi.setFocus()

if __name__ == "__main__":
    app = QApplication(sys.argv)
    pencere = YerelAsistanApp()
    pencere.show()
    sys.exit(app.exec_())